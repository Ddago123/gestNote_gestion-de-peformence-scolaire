"""Génère la présentation PowerPoint de GestNote."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Palette ──────────────────────────────────────────────
BLEU       = RGBColor(0x1e, 0x3a, 0x8a)   # bleu foncé primary
BLEU_CLAIR = RGBColor(0x37, 0x63, 0xca)   # bleu moyen
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR = RGBColor(0xF3, 0xF4, 0xF6)
GRIS_MOY   = RGBColor(0x6B, 0x72, 0x80)
VERT       = RGBColor(0x16, 0xa3, 0x4a)
ORANGE     = RGBColor(0xea, 0x58, 0x0c)
ROUGE      = RGBColor(0xdc, 0x26, 0x26)
NOIR       = RGBColor(0x1F, 0x29, 0x37)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout vide


# ─── Helpers ────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=BLEU, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = util.Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=BLANC,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   title=None, title_size=14, item_size=12,
                   title_color=BLEU, item_color=NOIR, bg=None, bullet="▸ "):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg)
    if title:
        add_text(slide, title, l+0.15, t+0.12, w-0.3, 0.35,
                 size=title_size, bold=True, color=title_color)
        offset = 0.45
    else:
        offset = 0.15
    for item in items:
        add_text(slide, bullet + item, l+0.15, t+offset, w-0.3, 0.32,
                 size=item_size, color=item_color)
        offset += 0.30
    return offset


def card(slide, l, t, w, h, title, items, title_color=BLEU,
         item_size=11.5, bg=GRIS_CLAIR, accent_color=BLEU):
    """Carte avec bande colorée à gauche."""
    add_rect(slide, l, t, 0.06, h, fill=accent_color)
    add_rect(slide, l+0.06, t, w-0.06, h, fill=bg)
    add_text(slide, title, l+0.2, t+0.1, w-0.25, 0.32,
             size=13, bold=True, color=title_color)
    for i, item in enumerate(items):
        add_text(slide, "• " + item, l+0.2, t+0.38+i*0.28, w-0.25, 0.28,
                 size=item_size, color=NOIR)


def badge(slide, text, l, t, w=1.2, h=0.35, fill=BLEU, text_color=BLANC, size=11):
    add_rect(slide, l, t, w, h, fill=fill)
    add_text(slide, text, l, t+0.03, w, h-0.05,
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Fond dégradé simulé : rectangle bleu foncé plein
add_rect(sl, 0, 0, 13.33, 7.5, fill=BLEU)
# Accent décoratif
add_rect(sl, 0, 5.8, 13.33, 1.7, fill=BLEU_CLAIR)
add_rect(sl, 0, 5.9, 13.33, 0.06, fill=BLANC)

# Logo / titre
add_text(sl, "GestNote", 1.0, 1.6, 8, 1.4,
         size=72, bold=True, color=BLANC)
add_text(sl, "Système de Gestion des Performances Scolaires",
         1.0, 3.1, 10, 0.7, size=22, color=RGBColor(0xBF, 0xDB, 0xFF))
add_text(sl, "Application web Flask · SQLite · JavaScript",
         1.0, 3.9, 10, 0.5, size=15, italic=True,
         color=RGBColor(0x93, 0xC5, 0xFD))

# Tags en bas
tags = [("Authentification sécurisée", VERT),
        ("Import CSV", ORANGE),
        ("Alertes pédagogiques", ROUGE),
        ("Export rapports", BLEU_CLAIR)]
x = 1.0
for label, col in tags:
    badge(sl, label, x, 6.15, 2.5, 0.42, fill=col)
    x += 2.7


# ════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE & OBJECTIFS
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Contexte & Objectifs", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "2", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# Problématique
add_rect(sl, 0.4, 1.3, 5.8, 2.8, fill=GRIS_CLAIR)
add_rect(sl, 0.4, 1.3, 0.08, 2.8, fill=ROUGE)
add_text(sl, "😟  Problématique", 0.6, 1.42, 5.4, 0.4,
         size=14, bold=True, color=ROUGE)
probs = [
    "Suivi des notes dispersé (papier, Excel)",
    "Aucune détection automatique des difficultés",
    "Pas de vision globale par classe/matière",
    "Absence d'alertes pédagogiques proactives",
    "Rapports manuels chronophages",
]
for i, p in enumerate(probs):
    add_text(sl, "• " + p, 0.6, 1.88+i*0.32, 5.4, 0.3, size=11.5, color=NOIR)

# Solution
add_rect(sl, 6.8, 1.3, 6.1, 2.8, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_rect(sl, 6.8, 1.3, 0.08, 2.8, fill=VERT)
add_text(sl, "✅  Solution GestNote", 7.0, 1.42, 5.7, 0.4,
         size=14, bold=True, color=VERT)
sols = [
    "Plateforme web centralisée et sécurisée",
    "Import CSV rapide des notes et élèves",
    "Calcul automatique des moyennes pondérées",
    "Alertes intelligentes (baisse, excellence…)",
    "Rapports et export CSV instantanés",
]
for i, s in enumerate(sols):
    add_text(sl, "• " + s, 7.0, 1.88+i*0.32, 5.7, 0.3, size=11.5, color=NOIR)

# Public cible
add_rect(sl, 0.4, 4.35, 12.5, 1.1, fill=RGBColor(0xFE, 0xF9, 0xEE))
add_rect(sl, 0.4, 4.35, 12.5, 0.06, fill=ORANGE)
add_text(sl, "👥  Public cible :", 0.6, 4.5, 3, 0.4,
         size=13, bold=True, color=ORANGE)
cibles = ["Directeurs d'établissements", "Enseignants", "Secrétaires pédagogiques"]
for i, c in enumerate(cibles):
    badge(sl, c, 3.2+i*3.3, 4.5, 2.9, 0.38, fill=ORANGE)

# Niveaux couverts
add_rect(sl, 0.4, 5.65, 12.5, 0.85, fill=GRIS_CLAIR)
add_text(sl, "🏫  Niveaux scolaires couverts :", 0.6, 5.75, 4, 0.4,
         size=12, bold=True, color=BLEU)
niveaux = "6ème · 5ème · 4ème · 3ème · 2nde · 1ère A/C/D · Terminale A/C/D"
add_text(sl, niveaux, 4.5, 5.75, 8.2, 0.4, size=12, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 3 — FONCTIONNALITÉS PRINCIPALES
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Fonctionnalités Principales", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "3", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

features = [
    ("🔐  Authentification", BLEU, [
        "Login sécurisé (session Flask)",
        "Credentials via variables d'env",
        "Protection de toutes les API",
        "Déconnexion automatique",
    ]),
    ("📋  Configuration", BLEU_CLAIR, [
        "Gestion années scolaires",
        "Création de classes par niveau",
        "Coefficients par matière",
        "Matières pré-configurées/niveau",
    ]),
    ("👨‍🎓  Gestion Élèves", VERT, [
        "Import CSV en masse",
        "Ajout individuel depuis l'UI",
        "Suppression avec cascade",
        "Recherche & navigation par classe",
    ]),
    ("📊  Import de Notes", ORANGE, [
        "Import CSV avec prévisualisation",
        "Détection des doublons",
        "Ajout manuel par formulaire",
        "Modification & suppression inline",
    ]),
    ("🔔  Alertes Pédagogiques", ROUGE, [
        "Baisse de performance détectée",
        "Notes basses consécutives",
        "Excellence maintenue",
        "Progression encourageante",
    ]),
    ("📈  Rapports & Stats", RGBColor(0x7C, 0x3A, 0xED), [
        "Moyenne pondérée par élève",
        "Écart-type par matière",
        "Graphique de distribution",
        "Export CSV compatible Excel",
    ]),
]

cols = [(0.3, 1.25), (4.55, 1.25), (8.8, 1.25),
        (0.3,  4.05), (4.55, 4.05), (8.8,  4.05)]
W, H = 3.9, 2.55

for (feat_title, accent, items), (cx, cy) in zip(features, cols):
    add_rect(sl, cx, cy, W, H, fill=GRIS_CLAIR)
    add_rect(sl, cx, cy, W, 0.07, fill=accent)
    add_text(sl, feat_title, cx+0.15, cy+0.15, W-0.2, 0.4,
             size=13, bold=True, color=accent)
    for i, item in enumerate(items):
        add_text(sl, "• " + item, cx+0.15, cy+0.55+i*0.44, W-0.2, 0.38,
                 size=10.5, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 4 — TABLEAU DE BORD
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Tableau de Bord", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "4", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# Zone alertes
add_rect(sl, 0.3, 1.3, 5.9, 2.75, fill=RGBColor(0xFE, 0xF2, 0xF2))
add_rect(sl, 0.3, 1.3, 5.9, 0.06, fill=ROUGE)
add_text(sl, "🔔  Alertes pédagogiques actives", 0.5, 1.4, 5.5, 0.4,
         size=13, bold=True, color=ROUGE)

alerte_exemples = [
    ("CRITIQUE", "Ali Koné — Mathématiques : 3 notes consécutives < 8", ROUGE, 5),
    ("ATTENTION", "Sara Diallo — Physique : baisse de performance", ORANGE, 3),
    ("SUCCÈS", "Marc Traoré — Français : excellence maintenue ≥17", VERT, 1),
]
for i, (typ, desc, col, sev) in enumerate(alerte_exemples):
    y = 1.88 + i * 0.66
    add_rect(sl, 0.45, y, 0.06, 0.52, fill=col)
    add_rect(sl, 0.51, y, 5.55, 0.52, fill=BLANC)
    badge(sl, typ, 0.6, y+0.08, 1.05, 0.28, fill=col, size=9)
    add_text(sl, desc, 1.75, y+0.06, 4.1, 0.4, size=10, color=NOIR)
    badge(sl, str(sev), 5.7, y+0.1, 0.28, 0.28, fill=col, size=10)

add_text(sl, "→ Bouton \"Résoudre\" individuel ou \"Résoudre toutes\"",
         0.5, 3.7, 5.7, 0.3, size=10, italic=True, color=GRIS_MOY)

# Zone stats matières
add_rect(sl, 6.55, 1.3, 6.45, 2.75, fill=GRIS_CLAIR)
add_rect(sl, 6.55, 1.3, 6.45, 0.06, fill=BLEU)
add_text(sl, "📊  Statistiques par matière", 6.75, 1.4, 6.0, 0.4,
         size=13, bold=True, color=BLEU)

headers = ["Matière", "Coef", "Moy.", "Éc.-type", "Min", "Max"]
col_x =   [6.7,       8.35,   8.95,   9.55,       10.35, 10.85]
col_w =   [1.55,      0.5,    0.55,   0.7,        0.45,  0.45]

for j, (h, cx) in enumerate(zip(headers, col_x)):
    add_text(sl, h, cx, 1.88, col_w[j], 0.3, size=9.5, bold=True, color=BLEU)

rows = [
    ("Mathématiques", "5", "13.45", "2.18", "8.0", "19.5"),
    ("Français",      "5", "11.20", "3.04", "5.5", "18.0"),
    ("Physique-Chim.","4", "10.80", "2.55", "6.0", "17.0"),
    ("Anglais",       "3", "14.10", "1.88", "10","19.0"),
]
for i, row in enumerate(rows):
    bg = BLANC if i % 2 == 0 else GRIS_CLAIR
    add_rect(sl, 6.65, 2.2+i*0.48, 6.2, 0.45, fill=bg)
    for j, (val, cx) in enumerate(zip(row, col_x)):
        col = GRIS_MOY
        if j == 2:  # Moyenne
            v = float(val)
            col = VERT if v >= 14 else (ORANGE if v >= 10 else ROUGE)
        add_text(sl, val, cx, 2.24+i*0.48, col_w[j], 0.35,
                 size=10, color=col, bold=(j==2))

# Zone liste élèves
add_rect(sl, 0.3, 4.25, 12.7, 2.9, fill=GRIS_CLAIR)
add_rect(sl, 0.3, 4.25, 12.7, 0.06, fill=BLEU_CLAIR)
add_text(sl, "👥  Liste des élèves — Moyennes & niveaux",
         0.5, 4.35, 9, 0.38, size=13, bold=True, color=BLEU)

eleves_ex = [
    ("EL001", "KOUASSI Jean",    "15.32", "Bon",   "0", VERT),
    ("EL002", "DIALLO Mariam",   "11.78", "Moyen", "1", ORANGE),
    ("EL003", "TRAORÉ Moussa",   "7.45",  "Faible","2", ROUGE),
    ("EL004", "BAMBA Aminata",   "13.90", "Moyen", "0", ORANGE),
]
ex = [0.5, 2.4, 5.2, 6.9, 8.8, 10.0]
hdrs2 = ["Matricule", "Nom Prénom", "Moyenne", "Niveau", "Alertes", ""]
for j, (h, cx) in enumerate(zip(hdrs2, ex)):
    add_text(sl, h, cx, 4.7, 1.6, 0.3, size=9.5, bold=True, color=BLEU)
for i, (mat, nom, moy, niv, al, col) in enumerate(eleves_ex):
    y = 5.06 + i * 0.42
    bg = BLANC if i % 2 == 0 else RGBColor(0xF9, 0xFA, 0xFB)
    add_rect(sl, 0.4, y, 12.5, 0.38, fill=bg)
    add_text(sl, mat, ex[0], y+0.04, 1.7, 0.3, size=10, color=GRIS_MOY)
    add_text(sl, nom, ex[1], y+0.04, 2.6, 0.3, size=10, color=NOIR, bold=True)
    add_text(sl, moy+"/20", ex[2], y+0.04, 1.4, 0.3, size=10, color=col, bold=True)
    badge(sl, niv, ex[3], y+0.05, 1.5, 0.28, fill=col, size=9)
    add_text(sl, al, ex[4]+0.5, y+0.04, 0.5, 0.3, size=10, color=ROUGE if al!="0" else VERT, bold=True)


# ════════════════════════════════════════════════════════
# SLIDE 5 — IMPORT & ALERTES
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Import de Notes & Système d'Alertes", 0.4, 0.2, 11, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "5", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# Flux import
add_rect(sl, 0.3, 1.25, 12.7, 0.06, fill=GRIS_CLAIR)
add_text(sl, "Flux d'import CSV", 0.3, 1.32, 5, 0.38, size=14, bold=True, color=BLEU)

steps = [
    ("1", "Sélectionner\nclasse & matière", BLEU),
    ("2", "Choisir le\nfichier CSV", BLEU_CLAIR),
    ("3", "Prévisualiser\n& valider", ORANGE),
    ("4", "Confirmer\nl'import", VERT),
    ("5", "Alertes\nautomatiques", ROUGE),
]
for i, (num, label, col) in enumerate(steps):
    x = 0.4 + i * 2.5
    add_rect(sl, x, 1.8, 2.1, 0.9, fill=col)
    add_text(sl, num, x+0.1, 1.82, 0.4, 0.4, size=20, bold=True, color=BLANC)
    add_text(sl, label, x+0.5, 1.85, 1.5, 0.8, size=10.5, color=BLANC, bold=True)
    if i < 4:
        add_text(sl, "→", x+2.12, 2.1, 0.3, 0.4, size=18, bold=True, color=BLEU)

# Détails format CSV
add_rect(sl, 0.3, 2.85, 5.9, 1.45, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_rect(sl, 0.3, 2.85, 0.07, 1.45, fill=BLEU)
add_text(sl, "📄  Format CSV attendu", 0.5, 2.95, 5.5, 0.35, size=12, bold=True, color=BLEU)
add_text(sl, "matricule, nom_eleve, note, date",
         0.5, 3.35, 5.5, 0.35, size=12, color=NOIR,
         italic=True)
add_text(sl, "EL001, KOUASSI Jean, 14.5, 2024-03-15\nEL002, DIALLO Mariam, 11.0, 2024-03-15",
         0.5, 3.7, 5.5, 0.5, size=10.5, color=GRIS_MOY, italic=True)

# Validations
add_rect(sl, 6.55, 2.85, 6.45, 1.45, fill=RGBColor(0xF0, 0xFD, 0xF4))
add_rect(sl, 6.55, 2.85, 0.07, 1.45, fill=VERT)
add_text(sl, "✅  Validations automatiques", 6.75, 2.95, 6, 0.35, size=12, bold=True, color=VERT)
validations = [
    "Note 0–20 vérifiée",
    "Matricule inconnu signalé",
    "Doublon (même élève, même date) ignoré",
    "Rapport : importées / ignorées / erreurs",
]
for i, v in enumerate(validations):
    add_text(sl, "• " + v, 6.75, 3.35+i*0.28, 6, 0.26, size=11, color=NOIR)

# Types d'alertes
add_rect(sl, 0.3, 4.5, 12.7, 0.05, fill=GRIS_CLAIR)
add_text(sl, "Types d'alertes générées automatiquement", 0.3, 4.6, 9, 0.38,
         size=14, bold=True, color=BLEU)

alertes_types = [
    ("Baisse de\nperformance", "Moyenne récente\n< ancienne – 2 pts", ROUGE, "Sévérité 3"),
    ("Notes basses\nconsécutives", "2 notes consécutives\nen dessous de 10", ORANGE, "Sévérité 3"),
    ("Difficulté\npersistante", "3 notes consécutives\nen dessous de 8", ROUGE, "Sévérité 5"),
    ("Excellence\nmaintenue", "Moyenne ≥ 17 sur\nles 5 dernières notes", VERT, "Sévérité 1"),
    ("Progression\nencourageante", "Hausse de 2 pts\nsur les 5 dernières", VERT, "Sévérité 2"),
]
for i, (titre, cond, col, sev) in enumerate(alertes_types):
    x = 0.35 + i * 2.57
    add_rect(sl, x, 5.05, 2.35, 1.9, fill=GRIS_CLAIR)
    add_rect(sl, x, 5.05, 2.35, 0.06, fill=col)
    add_text(sl, titre, x+0.1, 5.12, 2.1, 0.5, size=11, bold=True, color=col)
    add_text(sl, cond, x+0.1, 5.6, 2.1, 0.6, size=9.5, color=NOIR)
    badge(sl, sev, x+0.35, 6.55, 1.5, 0.3, fill=col, size=9)


# ════════════════════════════════════════════════════════
# SLIDE 6 — RAPPORTS & ANALYSE
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Rapports & Analyse", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "6", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# KPIs synthèse
kpis = [
    ("28", "Élèves", BLEU),
    ("12.4", "Moy. classe /20", VERT),
    ("11", "Bon niveau (≥14)", VERT),
    ("10", "Moyen (10–14)", ORANGE),
    ("7", "Faible (<10)", ROUGE),
    ("5", "Alertes actives", ROUGE),
]
for i, (val, label, col) in enumerate(kpis):
    x = 0.3 + i * 2.15
    add_rect(sl, x, 1.25, 2.0, 1.2, fill=col)
    add_text(sl, val, x+0.1, 1.3, 1.8, 0.65, size=32, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    add_text(sl, label, x+0.1, 1.88, 1.8, 0.45, size=10,
             color=BLANC, align=PP_ALIGN.CENTER)

# Graphique distribution (simulé)
add_rect(sl, 0.3, 2.65, 6.0, 3.3, fill=GRIS_CLAIR)
add_text(sl, "Distribution des niveaux (Chart.js)",
         0.5, 2.75, 5.5, 0.38, size=12, bold=True, color=BLEU)
bars = [("Bon (≥14)", 11, VERT), ("Moyen (10–14)", 10, ORANGE),
        ("Faible (<10)", 7, ROUGE)]
bar_max_h = 1.6
bar_max_val = 11
for i, (label, val, col) in enumerate(bars):
    x = 1.0 + i * 1.8
    h = bar_max_h * val / bar_max_val
    add_rect(sl, x, 4.95 - h, 1.2, h, fill=col)
    add_text(sl, str(val), x+0.35, 4.82-h, 0.5, 0.3, size=11, bold=True, color=col)
    add_text(sl, label, x, 5.0, 1.2, 0.38, size=8.5, color=GRIS_MOY,
             align=PP_ALIGN.CENTER)
add_rect(sl, 0.55, 4.96, 5.4, 0.04, fill=GRIS_MOY)

# Stats par matière
add_rect(sl, 6.55, 2.65, 6.45, 3.3, fill=GRIS_CLAIR)
add_text(sl, "Statistiques détaillées par matière",
         6.75, 2.75, 6.0, 0.38, size=12, bold=True, color=BLEU)

hdrs3 = ["Matière", "Coef", "Moy.", "Éc.type", "Min", "Max", "Notes"]
xcols3 = [6.65, 8.2, 8.75, 9.35, 10.1, 10.6, 11.15]
for j, (h, cx) in enumerate(zip(hdrs3, xcols3)):
    add_text(sl, h, cx, 3.15, 0.9, 0.28, size=9, bold=True, color=BLEU)

rows3 = [
    ("Mathématiques","5","13.45","2.18","8","19.5","56"),
    ("Français","5","11.20","3.04","5.5","18","56"),
    ("Physique-Chim.","4","10.80","2.55","6","17","42"),
    ("Anglais","3","14.10","1.88","10","19","28"),
    ("SVT","3","12.60","1.95","8.5","18","42"),
]
for i, row in enumerate(rows3):
    bg = BLANC if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFF)
    add_rect(sl, 6.6, 3.45+i*0.44, 6.3, 0.42, fill=bg)
    for j, (val, cx) in enumerate(zip(row, xcols3)):
        col2 = GRIS_MOY
        if j == 2:
            v = float(val)
            col2 = VERT if v >= 14 else (ORANGE if v >= 10 else ROUGE)
        add_text(sl, val, cx, 3.49+i*0.44, 0.85, 0.32,
                 size=9.5, color=col2, bold=(j==2))

# Export CSV
add_rect(sl, 0.3, 6.1, 12.7, 1.1, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_rect(sl, 0.3, 6.1, 12.7, 0.06, fill=BLEU)
add_text(sl, "📥  Export CSV  —", 0.5, 6.22, 2.5, 0.4, size=12, bold=True, color=BLEU)
add_text(sl, "Rapport complet : statistiques par matière + liste élèves avec moyennes  ·  Compatible Excel (BOM UTF-8)  ·  Téléchargement en 1 clic",
         3.1, 6.22, 9.8, 0.7, size=11, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE TECHNIQUE
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Architecture Technique", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "7", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# Backend
add_rect(sl, 0.3, 1.25, 3.9, 5.3, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_rect(sl, 0.3, 1.25, 3.9, 0.08, fill=BLEU)
add_text(sl, "🐍  Backend (Python/Flask)", 0.5, 1.38, 3.5, 0.38, size=13, bold=True, color=BLEU)
backend_items = [
    ("Flask Blueprints", "Routage modulaire"),
    ("SQLite3", "Base de données embarquée"),
    ("Flask Session", "Authentification"),
    ("Python csv/io", "Traitement des imports"),
    ("math (Python)", "Écart-type (SQLite sans STDEV)"),
    ("PRAGMA FK=ON", "Intégrité référentielle"),
    ("Connexion/appel", "Pattern simple & sûr"),
]
for i, (tech, desc) in enumerate(backend_items):
    y = 1.82 + i * 0.63
    add_rect(sl, 0.4, y, 3.7, 0.55, fill=BLANC)
    add_text(sl, tech, 0.55, y+0.04, 1.5, 0.28, size=10.5, bold=True, color=BLEU)
    add_text(sl, desc, 0.55, y+0.27, 3.4, 0.24, size=9.5, italic=True, color=GRIS_MOY)

# Flèche
add_text(sl, "⟺", 4.3, 3.5, 0.8, 0.5, size=22, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
add_text(sl, "REST\nJSON", 4.25, 4.1, 0.8, 0.5, size=9, color=GRIS_MOY, align=PP_ALIGN.CENTER)

# Frontend
add_rect(sl, 5.2, 1.25, 3.9, 5.3, fill=RGBColor(0xF0, 0xFD, 0xF4))
add_rect(sl, 5.2, 1.25, 3.9, 0.08, fill=VERT)
add_text(sl, "🌐  Frontend (Vanilla JS)", 5.4, 1.38, 3.5, 0.38, size=13, bold=True, color=VERT)
front_items = [
    ("SPA", "Single Page Application"),
    ("Fetch API", "Communication asynchrone"),
    ("Chart.js v4", "Graphiques (ligne & barres)"),
    ("escapeHtml()", "Protection XSS"),
    ("showToast()", "Notifications non-bloquantes"),
    ("setLoading()", "États de chargement boutons"),
    ("CSV export", "Blob + BOM UTF-8 (Excel)"),
]
for i, (tech, desc) in enumerate(front_items):
    y = 1.82 + i * 0.63
    add_rect(sl, 5.3, y, 3.7, 0.55, fill=BLANC)
    add_text(sl, tech, 5.45, y+0.04, 1.5, 0.28, size=10.5, bold=True, color=VERT)
    add_text(sl, desc, 5.45, y+0.27, 3.4, 0.24, size=9.5, italic=True, color=GRIS_MOY)

# Flèche
add_text(sl, "⟵", 9.2, 3.5, 0.7, 0.5, size=22, bold=True, color=BLEU, align=PP_ALIGN.CENTER)
add_text(sl, "GET\nSESSION", 9.15, 4.1, 0.8, 0.5, size=9, color=GRIS_MOY, align=PP_ALIGN.CENTER)

# Base de données
add_rect(sl, 10.0, 1.25, 3.0, 5.3, fill=RGBColor(0xFE, 0xF9, 0xEE))
add_rect(sl, 10.0, 1.25, 3.0, 0.08, fill=ORANGE)
add_text(sl, "🗄️  Base SQLite", 10.15, 1.38, 2.7, 0.38, size=13, bold=True, color=ORANGE)
tables = [
    "AnneesScolaires",
    "Classes",
    "Élèves",
    "Matières",
    "Notes",
    "Alertes",
]
for i, tbl in enumerate(tables):
    add_rect(sl, 10.1, 1.82+i*0.63, 2.8, 0.55, fill=BLANC)
    add_text(sl, "📋 " + tbl, 10.25, 1.9+i*0.63, 2.5, 0.35, size=11, bold=True, color=ORANGE)

# Performances
add_rect(sl, 0.3, 6.75, 12.7, 0.6, fill=GRIS_CLAIR)
add_text(sl, "⚡  Performances :", 0.5, 6.83, 2.2, 0.38, size=11, bold=True, color=BLEU)
perfs = [
    "3 requêtes fixes par chargement de classe (au lieu de 2N+3)",
    "Batch SQL : notes & alertes en 2 requêtes groupées",
    "Connexion légère par appel · Pas d'ORM overhead",
]
add_text(sl, "  ·  ".join(perfs), 2.8, 6.83, 10.1, 0.4, size=10, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 8 — EXPÉRIENCE UTILISATEUR
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Expérience Utilisateur", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "8", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

ux_items = [
    ("🔔  Toasts de notification",
     BLEU, ["Remplacent tous les alert() bloquants",
            "4 types : succès, erreur, avertissement, info",
            "Apparaissent 3.5s puis disparaissent seuls",
            "Position bas-droite, non-intrusifs"]),
    ("⚡  États de chargement",
     VERT, ["Spinner + bouton désactivé pendant requête",
            "6 boutons couverts (créer, importer, ajouter…)",
            "Évite les doubles clics et soumissions",
            "Retour à l'état initial après réponse"]),
    ("🖥️  Indicateurs de contenu",
     ORANGE, ["Spinner centré pendant chargement dashboard",
              "Spinner pendant chargement des rapports",
              "Remplacement instantané des placeholders",
              "Message d'erreur si la requête échoue"]),
    ("⌨️  Raccourcis clavier",
     RGBColor(0x7C, 0x3A, 0xED), ["Entrée sur champ \"Année\" → Créer",
                                   "Entrée sur champ \"Classe\" → Créer",
                                   "Entrée sur login/mdp → Se connecter",
                                   "Échap → Fermer la modal fiche élève"]),
    ("🔍  Navigation élèves",
     ROUGE, ["Sélecteur de classe pour lister tous les élèves",
             "Recherche textuelle (nom, prénom, matricule)",
             "Les deux modes partagent la zone de résultats",
             "Réinitialisation automatique entre les modes"]),
    ("✨  Polissage visuel",
     BLEU_CLAIR, ["Animation fluide d'ouverture de la modal",
                  "Bouton \"Détails\" classe avec toggle",
                  "Feedback ✓ vert sur sauvegarde coefficient",
                  "Classe CSS msg-warning ajoutée"]),
]

col_pos = [(0.3, 1.25), (4.55, 1.25), (8.8, 1.25),
           (0.3,  4.2),  (4.55, 4.2),  (8.8,  4.2)]
for (title, col, items), (cx, cy) in zip(ux_items, col_pos):
    W2, H2 = 3.9, 2.7
    add_rect(sl, cx, cy, W2, H2, fill=GRIS_CLAIR)
    add_rect(sl, cx, cy, W2, 0.07, fill=col)
    add_text(sl, title, cx+0.15, cy+0.13, W2-0.2, 0.38,
             size=12.5, bold=True, color=col)
    for i, item in enumerate(items):
        add_text(sl, "• " + item, cx+0.15, cy+0.55+i*0.48, W2-0.2, 0.42,
                 size=10.5, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 9 — SÉCURITÉ
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Sécurité", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "9", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

sec_items = [
    ("🔐  Authentification", BLEU, [
        "Session Flask avec secret_key configurable",
        "Identifiants via variables d'environnement",
        "hmac.compare_digest() contre les timing attacks",
        "before_request : toutes les API /api/* protégées",
        "401 → réaffichage automatique du formulaire de login",
    ]),
    ("🛡️  Protection XSS", ROUGE, [
        "Fonction escapeHtml() appliquée partout",
        "Noms d'élèves, alertes, classes, matricules",
        "Remplacement de innerHTML par du texte échappé",
        "Aucune donnée utilisateur insérée brute dans le DOM",
        "Couvre : recherche, dashboard, fiche élève, rapports",
    ]),
    ("🔒  Intégrité des données", VERT, [
        "PRAGMA foreign_keys = ON sur chaque connexion",
        "Requêtes paramétrées (aucun SQL injection possible)",
        "Contraintes UNIQUE en base (matricule, classe, nom)",
        "Détection de doublons côté serveur avant import",
        "Validation 0–20 en Python ET en SQLite (CHECK)",
    ]),
    ("⚙️  Configuration sécurisée", ORANGE, [
        "SECRET_KEY via $env:SECRET_KEY (jamais en dur)",
        "GESTNOTE_USER / GESTNOTE_PASS configurables",
        "Recommandation : SECRET_KEY aléatoire en prod",
        "Session permanente (durée de vie configurable)",
        "Pas de credentials dans le code source",
    ]),
]

positions = [(0.3, 1.25), (6.9, 1.25), (0.3, 4.05), (6.9, 4.05)]
for (title, col, items), (cx, cy) in zip(sec_items, positions):
    W3, H3 = 6.2, 2.6
    add_rect(sl, cx, cy, W3, H3, fill=GRIS_CLAIR)
    add_rect(sl, cx, cy, 0.08, H3, fill=col)
    add_text(sl, title, cx+0.22, cy+0.12, W3-0.3, 0.38,
             size=14, bold=True, color=col)
    for i, item in enumerate(items):
        add_text(sl, "• " + item, cx+0.22, cy+0.56+i*0.38, W3-0.3, 0.34,
                 size=11, color=NOIR)


# ════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION & ROADMAP
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.1, fill=BLEU)
add_text(sl, "Conclusion & Perspectives", 0.4, 0.2, 10, 0.7,
         size=28, bold=True, color=BLANC)
add_text(sl, "10", 12.5, 0.25, 0.6, 0.5, size=18, color=BLANC, align=PP_ALIGN.RIGHT)

# Bilan
add_rect(sl, 0.3, 1.25, 12.7, 2.0, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_rect(sl, 0.3, 1.25, 12.7, 0.07, fill=BLEU)
add_text(sl, "✅  Ce qui est livré", 0.5, 1.37, 5, 0.38, size=14, bold=True, color=BLEU)

bilan = [
    ("Configuration complète", "Années, classes, niveaux, coefficients"),
    ("Gestion des élèves", "Import CSV, ajout/suppression individuel, recherche, navigation"),
    ("Import de notes", "CSV avec preview, détection doublons, ajout manuel, édition inline"),
    ("Dashboard temps réel", "Alertes, stats matières, moyennes pondérées, fiche élève"),
    ("Alertes intelligentes", "5 types, 5 niveaux de sévérité, résolution depuis dashboard & fiche"),
    ("Rapports & Export", "Synthèse classe, distribution, écart-type, export CSV Excel"),
    ("Sécurité", "Auth session, XSS protégé, FK intégrité, requêtes paramétrées"),
    ("UX soignée", "Toasts, spinners, raccourcis clavier, animations, responsive"),
]
col1 = bilan[:4]
col2 = bilan[4:]
for i, (titre, desc) in enumerate(col1):
    add_text(sl, "✔ " + titre + " —", 0.5, 1.78+i*0.33, 3.5, 0.3, size=10.5, bold=True, color=BLEU)
    add_text(sl, desc, 3.5, 1.78+i*0.33, 3.0, 0.3, size=10.5, color=NOIR)
for i, (titre, desc) in enumerate(col2):
    add_text(sl, "✔ " + titre + " —", 6.9, 1.78+i*0.33, 3.5, 0.3, size=10.5, bold=True, color=BLEU)
    add_text(sl, desc, 10.0, 1.78+i*0.33, 3.0, 0.3, size=10.5, color=NOIR)

# Roadmap
add_rect(sl, 0.3, 3.45, 12.7, 0.06, fill=GRIS_CLAIR)
add_text(sl, "🗺️  Backlog — Prochaines évolutions", 0.3, 3.58, 8, 0.38,
         size=14, bold=True, color=BLEU)

roadmap = [
    ("#1", "Import tout-ou-rien\n→ lignes valides seules", ROUGE, "Robustesse"),
    ("#2", "try/catch manquants\n(supprimerClasse…)", ROUGE, "Robustesse"),
    ("#3", "Comparaison\ninter-classes", ORANGE, "Rapports"),
    ("#4", "Tri colonnes\npar clic d'en-tête", ORANGE, "UX"),
    ("#5", "Validation date\nnaissance (âge 8-25)", ORANGE, "Données"),
    ("#6", "Connexion BD\npar requête Flask (g)", RGBColor(0x7C,0x3A,0xED), "Perf."),
]
for i, (num, desc, col, cat) in enumerate(roadmap):
    x = 0.35 + i * 2.16
    add_rect(sl, x, 4.1, 2.0, 2.1, fill=GRIS_CLAIR)
    add_rect(sl, x, 4.1, 2.0, 0.06, fill=col)
    badge(sl, num, x+0.1, 4.18, 0.55, 0.32, fill=col, size=10)
    badge(sl, cat, x+0.72, 4.18, 1.15, 0.32, fill=RGBColor(0xFF,0xFF,0xFF),
          text_color=col, size=9)
    add_text(sl, desc, x+0.1, 4.57, 1.8, 0.75, size=10, color=NOIR)

# Stack technique recap
add_rect(sl, 0.3, 6.4, 12.7, 0.85, fill=BLEU)
techs = ["Python 3.12", "Flask 3.x", "SQLite 3", "Chart.js 4.4", "Vanilla JS ES6+", "python-pptx"]
add_text(sl, "Stack : ", 0.5, 6.55, 1.2, 0.4, size=12, bold=True, color=BLANC)
for i, t in enumerate(techs):
    badge(sl, t, 1.7+i*1.9, 6.55, 1.75, 0.38, fill=RGBColor(0x37,0x63,0xca), size=10)


# ── Sauvegarde ──────────────────────────────────────────
output = r"d:\Devs\Applis\gestion_de_prformance\GestNote_Presentation.pptx"
prs.save(output)
print("Fichier cree : " + output)
